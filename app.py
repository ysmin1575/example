from flask import Flask, render_template, request, send_file, redirect, url_for
from pptx import Presentation
import io
import os
import re
import requests

# -------------------------
# Groq AI 설정 (선택)
# -------------------------

try:
    from groq import Groq
    groq_key = os.environ.get("GROQ_API_KEY")

    if groq_key:
        client = Groq(api_key=groq_key)
    else:
        client = None
except:
    client = None

# -------------------------
# Flask 시작
# -------------------------

app = Flask(__name__)

# -------------------------
# 홈
# -------------------------

@app.route("/")
def home():
    return render_template("index.html")

# -------------------------
# PPT 생성 페이지
# -------------------------

@app.route("/generator")
def generator():
    return render_template("generator.html")

# -------------------------
# 미리보기 (AI 분석)
# -------------------------

@app.route("/preview", methods=["POST"])
def preview():

    title = request.form.get("title")
    name = request.form.get("name")
    content = request.form.get("content")

    # -------------------------
    # AI 사용
    # -------------------------

    if client:

        prompt = f"""
다음 내용을 PPT 슬라이드로 정리해줘.

조건
- 5~7 슬라이드
- 각 슬라이드는 한 문장
- 불필요한 기호 없이 문장만 출력

내용:
{content}
"""

        try:
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role":"user","content":prompt}]
            )

            result = response.choices[0].message.content.split("\n")

        except:
            result = content.split("\n")

    else:
        result = content.split("\n")

    # -------------------------
    # 텍스트 정리
    # -------------------------

    slides = []

    for s in result:

        s = s.strip()

        # 이상한 토큰 제거
        s = re.sub(r'_[A-Za-z0-9]+_', '', s)

        # 번호 제거
        s = re.sub(r'^\d+\.\s*', '', s)

        # 기호 제거
        s = re.sub(r'^[-•]\s*', '', s)

        # 마크다운 제거
        s = s.replace("**","").replace("`","")

        if s:
            slides.append(s)

    # 슬라이드 최소 보장
    if len(slides) < 3:
        slides = slides + ["추가 내용"]*(3-len(slides))

    return render_template(
        "preview.html",
        title=title,
        name=name,
        slides=slides
    )

# -------------------------
# PPT 생성
# -------------------------

@app.route("/generate", methods=["POST"])
def generate():

    title = request.form.get("title")
    name = request.form.get("name")
    slides = request.form.getlist("slides")

    prs = Presentation()

    # 제목 슬라이드
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = name

    # 내용 슬라이드
    for s in slides:

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = "내용"
        slide.placeholders[1].text = s

    # PPT 저장
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)

    return send_file(
        file_stream,
        as_attachment=True,
        download_name="slidr_ppt.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

# -------------------------
# 피드백 (Google Sheets)
# -------------------------

@app.route("/feedback", methods=["POST"])
def feedback():

    message = request.form.get("message")

    if message:

        try:

            url = "https://script.google.com/macros/s/AKfycbzYOD_D8KfUUKPB5r-_zBBtUCOMeR9SPfxOfqSufTvR814unuYT4pi5lWRwE7fc1IuImA/exec"

            requests.post(url, json={
                "message": message
            })

        except:
            pass

    return redirect(url_for("home"))

# -------------------------
# 서버 실행
# -------------------------

if __name__ == "__main__":

    port = int(os.environ.get("PORT", 10000))

    app.run(
        host="0.0.0.0",
        port=port
    )
